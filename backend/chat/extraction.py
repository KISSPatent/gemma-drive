"""
Text extraction for picked Drive files.
One function per source MIME family.
"""
import io
import csv
import pdfplumber
from googleapiclient.http import MediaIoBaseDownload

from pptx import Presentation
from docx import Document
from openpyxl import load_workbook


MAX_CHARS = 320_000  # ~80K tokens at ~4 chars/token

GOOGLE_DOC = "application/vnd.google-apps.document"
GOOGLE_SHEET = "application/vnd.google-apps.spreadsheet"
GOOGLE_SLIDES = "application/vnd.google-apps.presentation"
PDF = "application/pdf"
PLAIN_TEXT = "text/plain"
MARKDOWN = "text/markdown"
PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"



def _download_bytes(service, file_id: str) -> bytes:
    """Download non-Google-Workspace files."""
    request = service.files().get_media(fileId=file_id)
    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return buf.getvalue()


def _export_text(service, file_id: str, mime_type: str) -> str:
    """Export Google Workspace files as text."""
    request = service.files().export_media(fileId=file_id, mimeType=mime_type)
    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return buf.getvalue().decode("utf-8", errors="replace")


def extract(service, file_id: str, mime_type: str, name: str) -> tuple[str, bool]:
    """
    Returns (text, truncated).
    Raises on unsupported MIME or extraction failure.
    """
    if mime_type == PDF:
        data = _download_bytes(service, file_id)
        pages = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ""
                pages.append(t)
        text = "\n\n".join(pages)

    elif mime_type == GOOGLE_DOC:
        text = _export_text(service, file_id, "text/plain")

    elif mime_type == GOOGLE_SHEET:
        csv_text = _export_text(service, file_id, "text/csv")
        # Reformat as a markdown-ish table for readability
        reader = csv.reader(io.StringIO(csv_text))
        rows = list(reader)
        if rows:
            lines = []
            for row in rows[:1000]:  # cap rows
                lines.append(" | ".join(row))
            text = "\n".join(lines)
        else:
            text = ""

    elif mime_type == GOOGLE_SLIDES:
        text = _export_text(service, file_id, "text/plain")

    elif mime_type in (PLAIN_TEXT, MARKDOWN):
        text = _download_bytes(service, file_id).decode("utf-8", errors="replace")

    elif mime_type == PPTX:
        data = _download_bytes(service, file_id)
        prs = Presentation(io.BytesIO(data))
        slides_text = []
        for i, slide in enumerate(prs.slides, start=1):
            bits = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    bits.append(shape.text.strip())
            slides_text.append("\n".join(bits))
        text = "\n\n".join(slides_text)

    elif mime_type == DOCX:
        data = _download_bytes(service, file_id)
        doc = Document(io.BytesIO(data))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Also pull tables — common in patent docs, contracts, reports
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        text = "\n\n".join(parts)

    elif mime_type == XLSX:
        data = _download_bytes(service, file_id)
        wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                if row_count >= 1000:  # cap per sheet
                    rows.append(f"[... {ws.max_row - 1000} more rows truncated ...]")
                    break
                cells = ["" if c is None else str(c) for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
                    row_count += 1
            sheets_text.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
        wb.close()
        text = "\n\n".join(sheets_text)

    else:
        raise ValueError(f"Unsupported MIME type: {mime_type}")

    truncated = len(text) > MAX_CHARS
    if truncated:
        text = text[:MAX_CHARS] + f"\n\n[... truncated; full file is {len(text):,} chars ...]"

    return text, truncated


SUPPORTED_MIMES = {PDF, GOOGLE_DOC, GOOGLE_SHEET, GOOGLE_SLIDES, PLAIN_TEXT, MARKDOWN, PPTX, DOCX, XLSX}
